from __future__ import annotations
from typing import List, Tuple


class PptxExtractor:
    """Extract text from .pptx files, one slide per page."""

    def extract(self, path: str) -> List[Tuple[int, str]]:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("Install python-pptx: pip install 'docintel[office]'")

        prs = Presentation(path)
        pages: list[tuple[int, str]] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0)):
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if texts:
                pages.append((slide_num, "\n".join(texts)))

        return pages if pages else [(1, "")]
